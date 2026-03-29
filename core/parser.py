"""
core/parser.py
──────────────
Handles reading and cleaning PDF / PPTX files into plain text.
Supports text-based PDFs, table extraction, and OCR fallback.
"""

import re
from pathlib import Path
from typing import List, Optional

from config import COURSE_MATERIALS, CLEANED_TEXT_DIR


class SimpleDocParser:
    """Parse PDFs and PPTX files into clean .txt files."""

    SUPPORTED = {".pdf", ".pptx"}

    def __init__(
        self,
        input_folder: Path = COURSE_MATERIALS,
        output_folder: Path = CLEANED_TEXT_DIR,
    ):
        self.input_folder  = Path(input_folder)
        self.output_folder = Path(output_folder)
        self.output_folder.mkdir(exist_ok=True)
        self.processed: List[str] = []

    # ── Text Cleaning ─────────────────────────────────────────────────────────

    def clean_text(self, text: str) -> str:
        if not text:
            return ""
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r' +', ' ', text)
        text = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)
        text = re.sub(r'Page \d+ of \d+', '', text, flags=re.IGNORECASE)
        text = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f]', '', text)
        return text.strip()

    # ── Extractors ────────────────────────────────────────────────────────────

    def _extract_pdf(self, path: Path) -> str:
        import pdfplumber
        from pdf2image import convert_from_path
        import pytesseract

        parts = []
        try:
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        parts.append(f"--- Page {i+1} ---\n{page_text}")

                    for j, table in enumerate(page.extract_tables()):
                        if table:
                            rows = [
                                "\t".join(str(c) if c else "" for c in row)
                                for row in table
                            ]
                            parts.append(f"\n[Table {j+1} on Page {i+1}]\n" + "\n".join(rows))

            full = "\n\n".join(parts)

            # OCR fallback for scanned/image-only PDFs
            if len(full.strip()) < 100:
                print("   🔍 Sparse text — trying OCR...")
                images = convert_from_path(path)
                ocr_parts = [
                    f"--- Page {i+1} (OCR) ---\n{pytesseract.image_to_string(img)}"
                    for i, img in enumerate(images)
                    if pytesseract.image_to_string(img).strip()
                ]
                full = "\n\n".join(ocr_parts)

            return full

        except Exception as e:
            print(f"   ❌ PDF error: {e}")
            return ""

    def _extract_pptx(self, path: Path) -> str:
        from pptx import Presentation

        parts = []
        try:
            prs = Presentation(path)
            for i, slide in enumerate(prs.slides):
                slide_lines = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_lines.append(shape.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            slide_lines.append("\t".join(c.text for c in row.cells))
                if slide_lines:
                    parts.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_lines))
            return "\n\n".join(parts)

        except Exception as e:
            print(f"   ❌ PPTX error: {e}")
            return ""

    # ── Public API ────────────────────────────────────────────────────────────

    def process_file(self, file_path: Path) -> Optional[Path]:
        file_path = Path(file_path)

        if file_path.suffix.lower() not in self.SUPPORTED:
            print(f"⚠️  Skipping unsupported file: {file_path.name}")
            return None

        print(f"📖 Processing: {file_path.name}")

        raw = (
            self._extract_pdf(file_path)
            if file_path.suffix.lower() == ".pdf"
            else self._extract_pptx(file_path)
        )
        clean = self.clean_text(raw)

        if len(clean) < 50:
            print(f"   ⚠️  Insufficient text extracted")
            return None

        out = self.output_folder / f"{file_path.stem}.txt"
        out.write_text(
            f"Source: {file_path.name}\n{'='*60}\n\n{clean}",
            encoding="utf-8"
        )
        print(f"   ✅ Saved → {out.name}")
        self.processed.append(out.name)
        return out

    def process_folder(self, folder: Optional[Path] = None) -> None:
        folder = Path(folder or self.input_folder)
        files  = [f for f in folder.iterdir() if f.suffix.lower() in self.SUPPORTED]

        if not files:
            print(f"⚠️  No PDF or PPTX files found in {folder}/")
            return

        print(f"\n🗂  Processing {len(files)} file(s)...\n")
        for f in files:
            self.process_file(f)

        print(f"\n✨ Done — {len(self.processed)} file(s) parsed → {self.output_folder}/")
